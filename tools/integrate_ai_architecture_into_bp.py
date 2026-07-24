from __future__ import annotations

import os
import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
TARGET = ROOT / "平台资料" / "融资PPT" / "小熊AI - AI研发全链路平台 · 天使轮路演BP.pptx"
TEAM_SOURCE = ROOT / "小熊AI-融资路演BP-v2 - 副本.pptx"
OUTPUT = Path(
    os.environ.get(
        "BP_OUTPUT",
        ROOT / "平台资料" / "融资PPT" / "小熊AI - AI研发全链路平台 · 天使轮路演BP-含架构与团队页.pptx",
    )
)

EMU_PER_CM = 360000

BG = "0B1220"
PANEL = "151D2E"
PANEL_2 = "101827"
BORDER = "1E3050"
CYAN = "00D4FF"
GOLD = "FFB800"
PURPLE = "A855F7"
GREEN = "00C853"
RED = "FF4D4F"
TEXT = "FFFFFF"
MUTED = "94A3B8"
FOOTER = "64748B"
FONT = "MiSans"
DISPLAY_NUMBER_FONT = "Liter"


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def cm(value: float) -> int:
    return int(value * EMU_PER_CM)


def add_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, cm(x), cm(y), cm(w), cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1)
    return shp


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, valign=None):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multiline(slide, lines, x, y, w, h, size=13, color=MUTED, bullet=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"· {line}" if bullet else line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
    return box


def add_top_bar(slide, color=CYAN):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, cm(45.16), cm(0.105))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.color.rgb = rgb(color)


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, cm(45.16), cm(25.4))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(BG)
    bg.line.color.rgb = rgb(BG)


def add_title(slide, title, subtitle, section):
    add_background(slide)
    add_top_bar(slide, CYAN)
    add_text(slide, title, 2.12, 1.06, 31.75, 1.25, size=36, bold=True)
    add_text(slide, subtitle, 2.12, 2.45, 31.2, 0.48, size=14, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(2.12), cm(3.0), cm(3.6), cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(CYAN)
    line.line.color.rgb = rgb(CYAN)
    add_text(slide, section, 36.69, 23.99, 6.35, 0.85, size=14, color=FOOTER, align=PP_ALIGN.RIGHT)


def add_connector(slide, x1, y1, x2, y2, color="#607D95", dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, cm(x1), cm(y1), cm(x2), cm(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(1.5)
    if dashed:
        conn.line.dash_style = 4
    dx = x2 - x1
    dy = y2 - y1
    if abs(dx) >= abs(dy):
        rotation = 0 if dx >= 0 else 180
    else:
        rotation = 90 if dy >= 0 else 270
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cm(x2 - 0.13), cm(y2 - 0.13), cm(0.26), cm(0.26))
    tri.rotation = rotation
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(color)
    tri.line.color.rgb = rgb(color)
    return conn


def add_pill(slide, text, x, y, w, fill, color=TEXT, size=11):
    add_rect(slide, x, y, w, 0.58, fill=fill, line=BORDER, radius=True)
    add_text(slide, text, x + 0.08, y + 0.14, w - 0.16, 0.3, size=size, color=color, align=PP_ALIGN.CENTER)


def replace_text(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def add_product_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "AI平台产品架构图",
        "AI产品承接生成结果，AI会话作为执行中枢，AI开发完成真实研发闭环，AI军火库提供可复用能力资产。",
        "04A / 产品架构补充",
    )

    # Center session.
    add_rect(slide, 13.65, 4.0, 17.8, 5.3, fill=PANEL, line=CYAN)
    add_text(slide, "AI会话", 13.65, 4.45, 17.8, 0.65, size=26, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "统一交互入口 / 上下文编排 / 能力执行 / 工作区治理", 14.45, 5.35, 16.2, 0.45, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["会话消息", "文件上下文", "模型环境", "工作区资产"]):
        x = 15.0 + (i % 2) * 7.4
        y = 6.15 + (i // 2) * 1.3
        add_pill(slide, item, x, y, 5.5, fill="1A3A5A", color=TEXT, size=11)

    # Left product.
    add_rect(slide, 2.1, 5.0, 9.1, 6.2, fill=PANEL, line=GREEN)
    add_text(slide, "AI产品", 2.1, 5.45, 9.1, 0.55, size=23, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_multiline(
        slide,
        ["产品创建 / 产品列表", "项目工作台 / 产品预览", "生成服务 URL 识别", "会话结果转产品资产"],
        2.85,
        6.55,
        7.6,
        3.55,
        size=13,
        color=TEXT,
        bullet=True,
    )

    # Right workshop.
    add_rect(slide, 33.95, 5.0, 9.1, 6.2, fill=PANEL, line=GOLD)
    add_text(slide, "AI军火库", 33.95, 5.45, 9.1, 0.55, size=23, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_multiline(
        slide,
        ["Skill / Prompt / Agent / MCP", "大模型 / Git / 测试工具", "能力审核 / 版本 / 密钥", "添加到本地 Codex Runtime"],
        34.7,
        6.55,
        7.6,
        3.55,
        size=13,
        color=TEXT,
        bullet=True,
    )

    # Bottom dev.
    add_rect(slide, 8.0, 14.05, 29.2, 4.2, fill=PANEL, line=PURPLE)
    add_text(slide, "AI开发", 8.0, 14.45, 29.2, 0.55, size=23, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "项目配置、任务执行、缺陷同步、构建部署闭环", 9.0, 15.18, 27.2, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    dev_items = ["开发项目配置", "AI开发任务", "执行日志与缺陷", "构建制品部署"]
    for i, item in enumerate(dev_items):
        add_pill(slide, item, 9.8 + i * 6.6, 16.25, 5.2, fill="2A1A3A", color=TEXT, size=11)

    add_connector(slide, 11.2, 7.8, 13.65, 6.3, color=GREEN)
    add_connector(slide, 31.45, 6.3, 33.95, 7.8, color=GOLD)
    add_connector(slide, 22.55, 9.3, 22.55, 14.05, color=PURPLE)
    add_connector(slide, 15.0, 14.05, 10.8, 11.2, color=MUTED, dashed=True)
    add_connector(slide, 30.2, 14.05, 35.2, 11.2, color=MUTED, dashed=True)

    add_rect(slide, 3.0, 20.1, 39.0, 1.0, fill="0A1A10", line=GREEN)
    add_text(
        slide,
        "典型闭环：AI军火库选择能力 → AI会话组织上下文并执行 → AI开发生成/验证/部署 → AI产品沉淀可访问应用与交付结果",
        3.3,
        20.42,
        38.4,
        0.35,
        size=12,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    return slide


def add_tech_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "AI平台技术架构图",
        "Vue 前端承载产品入口，Django REST 提供会话/能力/开发/部署 API，Celery 与 Codex Runtime 执行真实开发链路。",
        "04B / 技术架构补充",
    )

    layers = [
        ("前端表现层（Vue / Vite / Element Plus）", CYAN, 3.7, [
            ("AIProductCreate.vue", "AI产品"),
            ("CodexChat / ConversationWorkspace", "AI会话"),
            ("AI开发任务 / 工作区面板", "AI开发"),
            ("AIWorkshopManager.vue", "AI军火库"),
        ]),
        ("后端 API 与领域服务层（Django REST Framework）", GREEN, 8.2, [
            ("apps.assistant", "sessions / chat / files / capabilities"),
            ("apps.ai_development", "configs / tasks / defects"),
            ("apps.deployments", "targets / artifacts / executions / rollback"),
            ("权限 / 审核 / 配置", "PermissionItem / review / secrets"),
        ]),
        ("执行运行时层", PURPLE, 13.05, [
            ("Celery Worker", "异步任务编排"),
            ("Codex CLI Runtime", "会话驱动开发"),
            ("AI Tool Controller", "模型与工具调用"),
            ("Docker / Git", "构建、测试、提交"),
        ]),
        ("数据与外部依赖层", GOLD, 17.55, [
            ("MySQL / Redis", "会话、消息、资产、任务、队列"),
            ("文件与工作区存储", "上传文件、代码资产、构建产物"),
            ("外部服务", "LLM / Codex / Git / Docker / 部署目标"),
        ]),
    ]

    for title, color, y, cards in layers:
        add_rect(slide, 2.1, y, 40.95, 3.1 if len(cards) == 4 else 2.55, fill=PANEL, line=BORDER)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(2.1), cm(y), cm(0.12), cm(3.1 if len(cards) == 4 else 2.55))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(color)
        stripe.line.color.rgb = rgb(color)
        add_text(slide, title, 2.65, y + 0.28, 20.0, 0.45, size=16, color=color, bold=True)
        n = len(cards)
        card_w = 9.0 if n == 4 else 12.45
        gap = 0.82 if n == 4 else 0.95
        start_x = 3.0
        for i, (name, desc) in enumerate(cards):
            x = start_x + i * (card_w + gap)
            add_rect(slide, x, y + 1.02, card_w, 1.35, fill=PANEL_2, line=BORDER)
            add_text(slide, name, x + 0.25, y + 1.22, card_w - 0.5, 0.35, size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, desc, x + 0.25, y + 1.77, card_w - 0.5, 0.3, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    add_connector(slide, 22.55, 6.8, 22.55, 8.2, color=MUTED)
    add_connector(slide, 22.55, 11.3, 22.55, 13.05, color=MUTED)
    add_connector(slide, 22.55, 16.15, 22.55, 17.55, color=MUTED)

    add_connector(slide, 16.2, 6.2, 8.1, 9.2, color=CYAN, dashed=True)
    add_connector(slide, 25.0, 6.2, 17.9, 9.2, color=CYAN, dashed=True)
    add_connector(slide, 34.5, 6.2, 35.0, 9.2, color=CYAN, dashed=True)
    add_connector(slide, 17.9, 11.3, 16.7, 14.05, color=GREEN, dashed=True)
    add_connector(slide, 27.8, 11.3, 36.0, 14.05, color=GREEN, dashed=True)
    return slide


def add_business_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "商业模式 · 订阅 + 用量 + 私有化",
        "基础订阅形成可持续 ARR，用量计费捕获高频 AI 执行价值，企业私有化打开高客单价政企市场。",
        "10 / 商业模式",
    )

    cards = [
        (
            "SaaS 订阅",
            CYAN,
            "10-50万 / 年",
            ["按团队规模、项目数、AI会话额度分级", "覆盖新建研发团队与中小型企业研发组织", "标准化交付，形成稳定 ARR"],
        ),
        (
            "企业私有化部署",
            GOLD,
            "50-300万 / 年",
            ["面向金融、国企、医疗、政务等高合规客户", "支持内网部署、模型隔离、日志审计", "高客单价，绑定企业 AI 研发基础设施"],
        ),
        (
            "AI 能力用量计费",
            PURPLE,
            "按执行量计费",
            ["Token、AI任务、自动化测试、构建部署次数", "与客户真实使用强相关", "高频研发活动带来持续增量收入"],
        ),
        (
            "专业服务与实施",
            GREEN,
            "10-100万 / 项目",
            ["需求基线、研发流程、质量体系接入", "Git / Docker / CI-CD / 钉钉 / 飞书 / 企微集成", "降低首批政企客户落地阻力"],
        ),
        (
            "生态能力市场",
            RED,
            "订阅 / 抽佣",
            ["Skill、Prompt、Agent、MCP 能力包", "官方能力 + 行业能力 + 第三方能力", "从工具收费升级为能力资产分发平台"],
        ),
    ]

    positions = [
        (2.1, 3.55, 12.65, 6.45),
        (16.25, 3.55, 12.65, 6.45),
        (30.4, 3.55, 12.65, 6.45),
        (2.1, 11.0, 19.8, 5.65),
        (23.25, 11.0, 19.8, 5.65),
    ]

    for idx, ((title, color, price, lines), (x, y, w, h)) in enumerate(zip(cards, positions), start=1):
        add_rect(slide, x, y, w, h, fill=PANEL, line=color)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(x), cm(y), cm(w), cm(0.12))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(color)
        stripe.line.color.rgb = rgb(color)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(x + 0.55), cm(y + 0.58), cm(1.05), cm(1.05))
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(color)
        badge.line.color.rgb = rgb(color)
        add_text(slide, str(idx), x + 0.55, y + 0.79, 1.05, 0.3, size=13, color=BG, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, title, x + 1.85, y + 0.55, w - 2.35, 0.45, size=18, color=color, bold=True)
        add_text(slide, price, x + 1.85, y + 1.2, w - 2.35, 0.4, size=13, color=TEXT, bold=True)
        add_multiline(slide, lines, x + 0.85, y + 2.05, w - 1.45, h - 2.55, size=11.5, color=TEXT, bullet=True)

    add_rect(slide, 2.1, 18.25, 40.95, 1.45, fill="0A1A10", line=GREEN)
    add_text(
        slide,
        "收入结构 = 年费订阅 ARR + AI执行用量 + 私有化部署 + 实施集成 + 能力市场分成",
        2.55,
        18.68,
        40.05,
        0.4,
        size=16,
        color=TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        "商业逻辑：先用标准化 SaaS 跑通标杆客户，再用治理、审计、私有化与生态能力市场提升客单价和续费粘性。",
        3.1,
        20.35,
        38.95,
        0.45,
        size=12.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    return slide


def add_team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "团队",
        "创始人具备 16 年互联网与 AI 研发质量经验，平台本身由一个人纯用 AI 完成，是能力验证也是产品背书。",
        "12 / 团队",
    )

    add_rect(slide, 2.1, 3.1, 16.2, 15.0, fill=PANEL, line=BORDER)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(2.1), cm(3.1), cm(16.2), cm(0.14))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(GOLD)
    stripe.line.color.rgb = rgb(GOLD)

    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, cm(6.95), cm(4.25), cm(6.3), cm(6.3))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = rgb("1A2A4A")
    avatar.line.color.rgb = rgb(GOLD)
    avatar.line.width = Pt(2)
    add_text(slide, "👤", 8.47, 5.9, 3.25, 1.1, size=48, align=PP_ALIGN.CENTER)
    add_text(slide, "冯绍文", 2.7, 11.15, 15.0, 0.65, size=25, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "创始人 & CEO", 2.7, 12.05, 15.0, 0.45, size=15, color=MUTED, align=PP_ALIGN.CENTER)

    tags = [
        ("高级测试经理", "1A3A5A"),
        ("16年互联网+AI经验", "1A3A2A"),
        ("深圳国企・生物基因龙头", "2A1A3A"),
        ("带过30+人测试团队", "3A1A1A"),
        ("攻坚3亿+体量项目", "2A2A0A"),
    ]
    positions = [(3.0, 13.25, 6.2), (10.0, 13.25, 6.2), (3.0, 14.25, 6.2), (10.0, 14.25, 6.2), (6.5, 15.25, 6.2)]
    for (txt, fill), (x, y, w) in zip(tags, positions):
        add_pill(slide, txt, x, y, w, fill=fill, size=10.5)

    add_rect(slide, 20.0, 3.1, 22.9, 7.4, fill="0A1A10", line=BORDER)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(20.0), cm(3.1), cm(22.9), cm(0.14))
    top.fill.solid()
    top.fill.fore_color.rgb = rgb(GREEN)
    top.line.color.rgb = rgb(GREEN)
    add_text(slide, "“一个人，纯用AI，", 20.9, 4.0, 20.9, 0.85, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "构建了一个10人团队", 20.9, 5.25, 20.9, 0.85, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "才能构建的平台”", 20.9, 6.5, 20.9, 0.85, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_pill(slide, "100% AI开发 · 50+亿Token", 27.7, 8.45, 7.2, fill="0A2A1A", color=GREEN, size=12)

    add_rect(slide, 20.0, 11.45, 22.9, 5.45, fill=PANEL, line=BORDER)
    add_multiline(
        slide,
        [
            "这本身，就是本平台能力最真实的背书",
            "",
            "我们不只是在卖一款工具",
            "我们是在证明一种新的研发范式",
        ],
        21.0,
        12.25,
        20.8,
        3.4,
        size=22,
        color=TEXT,
        bullet=False,
    )

    add_rect(slide, 2.1, 19.35, 40.8, 1.05, fill="151D2E", line=RED)
    add_text(slide, "创始人即平台首个重度用户：用 AI 会话、AI军火库与 AI开发闭环完成平台自身迭代", 2.45, 19.7, 40.1, 0.35, size=13, color=TEXT, align=PP_ALIGN.CENTER)
    return slide


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    slide = slides[old_index]
    sld_id_lst.remove(slide)
    sld_id_lst.insert(new_index, slide)


def add_and_reorder(prs: Presentation) -> None:
    original_count = len(prs.slides)
    add_product_architecture_slide(prs)
    add_tech_architecture_slide(prs)
    add_business_model_slide(prs)
    add_team_slide(prs)
    # Move product and tech pages after current slide 5.
    move_slide(prs, original_count, 5)
    move_slide(prs, original_count + 1, 6)
    # Move business model page after market size and before commercialization path.
    move_slide(prs, original_count + 2, 12)
    # Move team page before financing page.
    move_slide(prs, original_count + 3, 14)
    replace_text(prs.slides[13], "10 / 商业化路径", "11 / 商业化路径")
    replace_text(prs.slides[15], "11 / 融资计划", "13 / 融资计划")


def is_footer_text(text: str) -> bool:
    return bool(re.fullmatch(r"\d{2}[A-Z]? / .+", text.strip()))


def normalize_text_frame(shape) -> None:
    if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
        return

    text = shape.text.strip()
    is_footer = is_footer_text(text)

    tf = shape.text_frame
    # Keep text boxes visually tight and consistent with the existing deck style.
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for paragraph in tf.paragraphs:
        if is_footer:
            paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            if not run.text:
                continue
            current_font = run.font.name
            if current_font != DISPLAY_NUMBER_FONT:
                run.font.name = FONT
            if is_footer:
                run.font.size = Pt(14)
                run.font.bold = False
                run.font.color.rgb = rgb(FOOTER)
        if is_footer:
            shape.left = cm(35.98)
            shape.top = cm(23.99)
            shape.width = cm(7.06)
            shape.height = cm(0.85)


def normalize_presentation_format(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            normalize_text_frame(shape)


def main():
    prs = Presentation(str(TARGET))
    add_and_reorder(prs)
    normalize_presentation_format(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)
    print("slides", len(prs.slides))


if __name__ == "__main__":
    main()
