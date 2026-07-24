import os
import json
import time
import uuid
import asyncio
import logging
import shutil
import subprocess
import tarfile
import tempfile
import zipfile
from io import BytesIO
import xml.etree.ElementTree as ET
from typing import List, Dict, Any, Optional
from datetime import datetime

try:
    from PyPDF2 import PdfReader
except ImportError:
    from PyPDF2 import PdfFileReader as PdfReader
    
try:
    import docx
except ImportError:
    docx = None
from django.conf import settings
from django.core.files.storage import default_storage

from .models import RequirementDocument, RequirementAnalysis, BusinessRequirement, GeneratedTestCase, AnalysisTask

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Document text extraction service."""

    ARCHIVE_EXTENSIONS = ('.tar.gz', '.tgz', '.zip', '.tar', '.rar')
    DOCUMENT_EXTENSIONS = (
        '.pdf', '.doc', '.docx', '.txt', '.md', '.xmind',
        '.png', '.jpg', '.jpeg', '.xls', '.xlsx', '.ppt', '.pptx',
    )
    MAX_ARCHIVE_FILES = 60
    MAX_ARCHIVE_MEMBER_BYTES = 30 * 1024 * 1024
    MAX_ARCHIVE_TOTAL_BYTES = 120 * 1024 * 1024
    MAX_EXTRACTED_TEXT_CHARS = 200_000

    @classmethod
    def get_file_extension(cls, filename: str) -> str:
        lower_name = str(filename or '').lower()
        for extension in cls.ARCHIVE_EXTENSIONS:
            if lower_name.endswith(extension):
                return extension
        return os.path.splitext(lower_name)[1]

    @classmethod
    def is_supported_document_file(cls, filename: str, allow_archives: bool = True) -> bool:
        extension = cls.get_file_extension(filename)
        return extension in cls.DOCUMENT_EXTENSIONS or (allow_archives and extension in cls.ARCHIVE_EXTENSIONS)

    @classmethod
    def truncate_text(cls, text: str) -> str:
        if len(text) <= cls.MAX_EXTRACTED_TEXT_CHARS:
            return text
        return f"{text[:cls.MAX_EXTRACTED_TEXT_CHARS]}\n...[truncated]"

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        try:
            lines = []
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text() or ''
                    if page_text.strip():
                        lines.append(page_text.strip())
            return '\n'.join(lines).strip()
        except Exception as e:
            logger.error(f"PDF文本提取失败: {e}")
            return f"PDF文本提取失败: {str(e)}"

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        try:
            if docx is None:
                return 'Word文档文本提取失败: python-docx 未安装'

            document = docx.Document(file_path)
            lines = []
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines.append(text)
            for table in document.tables:
                for row in table.rows:
                    values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if values:
                        lines.append('\t'.join(values))
            return '\n'.join(lines).strip()
        except Exception as e:
            logger.error(f"Word文档文本提取失败: {e}")
            return f"Word文档文本提取失败: {str(e)}"

    @staticmethod
    def extract_text_from_legacy_office(file_path: str, file_type: str = 'Office') -> str:
        """Best-effort extraction for legacy binary .doc/.ppt files."""
        command = None
        if file_type == 'Word' and shutil.which('antiword'):
            command = ['antiword', file_path]
        elif shutil.which('strings'):
            command = ['strings', file_path]

        if command:
            try:
                result = subprocess.run(command, capture_output=True, text=True, timeout=30, check=False)
                output = (result.stdout or '').strip()
                if output:
                    return output
            except Exception as e:
                logger.warning(f"{file_type} legacy command extraction failed: {e}")

        return DocumentProcessor.extract_text_from_binary(file_path, file_type=file_type)

    @staticmethod
    def extract_text_from_binary(file_path: str, file_type: str = 'Binary') -> str:
        try:
            with open(file_path, 'rb') as file:
                raw = file.read(2 * 1024 * 1024)
            for encoding in ('utf-8', 'gbk', 'utf-16', 'latin-1'):
                try:
                    text = raw.decode(encoding, errors='ignore')
                    cleaned = ''.join(char if char.isprintable() or char in '\r\n\t' else ' ' for char in text)
                    cleaned = '\n'.join(line.strip() for line in cleaned.splitlines() if line.strip())
                    if cleaned:
                        return cleaned
                except Exception:
                    continue
            return f"{file_type}文本提取失败: 未识别到可读文本"
        except Exception as e:
            logger.error(f"{file_type}文本提取失败: {e}")
            return f"{file_type}文本提取失败: {str(e)}"

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        for encoding in ('utf-8-sig', 'utf-8', 'gbk', 'utf-16', 'latin-1'):
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read().strip()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"文本文件读取失败: {e}")
                return f"文本文件读取失败: {str(e)}"

        try:
            with open(file_path, 'rb') as file:
                return file.read().decode('utf-8', errors='ignore').strip()
        except Exception as e:
            logger.error(f"文本文件读取失败: {e}")
            return f"文本文件读取失败: {str(e)}"

    @staticmethod
    def extract_text_from_xmind(file_path: str) -> str:
        try:
            with zipfile.ZipFile(file_path, 'r') as archive:
                names = set(archive.namelist())
                if 'content.json' in names:
                    payload = json.loads(archive.read('content.json').decode('utf-8'))
                    lines = DocumentProcessor._extract_xmind_json_lines(payload)
                elif 'content.xml' in names:
                    payload = archive.read('content.xml').decode('utf-8')
                    lines = DocumentProcessor._extract_xmind_xml_lines(payload)
                else:
                    lines = []
                    for name in archive.namelist():
                        if not (name.endswith('/content.json') or name.endswith('/content.xml')):
                            continue
                        raw = archive.read(name).decode('utf-8')
                        if name.endswith('.json'):
                            lines.extend(DocumentProcessor._extract_xmind_json_lines(json.loads(raw)))
                        else:
                            lines.extend(DocumentProcessor._extract_xmind_xml_lines(raw))

            normalized = []
            seen = set()
            for line in lines:
                text = str(line or '').strip()
                if text and text not in seen:
                    normalized.append(text)
                    seen.add(text)
            return '\n'.join(normalized)
        except Exception as e:
            logger.error(f"XMind文本提取失败: {e}")
            return f"XMind文本提取失败: {str(e)}"

    @staticmethod
    def _extract_xmind_json_lines(payload) -> list[str]:
        lines = []

        def walk(node, depth=0):
            if isinstance(node, dict):
                title = str(node.get('title') or '').strip()
                if title:
                    lines.append(f"{'  ' * depth}- {title}")

                notes = node.get('notes') or {}
                plain = notes.get('plain') if isinstance(notes, dict) else None
                content = plain.get('content') if isinstance(plain, dict) else None
                if content:
                    lines.append(f"{'  ' * (depth + 1)}备注: {str(content).strip()}")

                children = node.get('children')
                if isinstance(children, dict):
                    for child_group in children.values():
                        if isinstance(child_group, list):
                            for child in child_group:
                                walk(child, depth + 1)
                        else:
                            walk(child_group, depth + 1)

                for key in ('rootTopic', 'topic', 'topics'):
                    value = node.get(key)
                    if value is not None:
                        walk(value, depth)
            elif isinstance(node, list):
                for item in node:
                    walk(item, depth)

        walk(payload)
        return lines

    @staticmethod
    def _extract_xmind_xml_lines(payload: str) -> list[str]:
        lines = []
        root = ET.fromstring(payload)
        for element in root.iter():
            tag = element.tag.rsplit('}', 1)[-1]
            if tag in {'title', 'plain'}:
                text = ''.join(element.itertext()).strip()
                if text:
                    lines.append(f'- {text}')
        return lines

    @staticmethod
    def extract_text_from_excel(file_path: str) -> str:
        extension = DocumentProcessor.get_file_extension(file_path)
        try:
            if extension == '.xls':
                try:
                    import xlrd
                except ImportError:
                    return 'Excel文本提取失败: xlrd 未安装，无法解析 .xls 文件'

                workbook = xlrd.open_workbook(file_path)
                lines = []
                for sheet in workbook.sheets():
                    lines.append(f"[Sheet] {sheet.name}")
                    for row_index in range(sheet.nrows):
                        values = [
                            str(sheet.cell_value(row_index, col_index)).strip()
                            for col_index in range(sheet.ncols)
                            if str(sheet.cell_value(row_index, col_index)).strip()
                        ]
                        if values:
                            lines.append('\t'.join(values))
                return '\n'.join(lines).strip()

            try:
                import openpyxl
            except ImportError:
                return 'Excel文本提取失败: openpyxl 未安装'

            workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            lines = []
            for worksheet in workbook.worksheets:
                lines.append(f"[Sheet] {worksheet.title}")
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        lines.append('\t'.join(values))
            workbook.close()
            return '\n'.join(lines).strip()
        except Exception as e:
            logger.error(f"Excel文本提取失败: {e}")
            return f"Excel文本提取失败: {str(e)}"

    @staticmethod
    def extract_text_from_ppt(file_path: str) -> str:
        extension = DocumentProcessor.get_file_extension(file_path)
        if extension == '.ppt':
            return DocumentProcessor.extract_text_from_legacy_office(file_path, file_type='PPT')

        try:
            try:
                from pptx import Presentation
            except ImportError:
                return 'PPT文本提取失败: python-pptx 未安装'

            presentation = Presentation(file_path)
            lines = []

            def collect_shape_text(shape):
                if getattr(shape, 'has_text_frame', False):
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if values:
                            lines.append('\t'.join(values))
                if hasattr(shape, 'shapes'):
                    for child_shape in shape.shapes:
                        collect_shape_text(child_shape)

            for index, slide in enumerate(presentation.slides, start=1):
                lines.append(f"[Slide {index}]")
                for shape in slide.shapes:
                    collect_shape_text(shape)
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        lines.append(f"[Notes]\n{notes}")
                except Exception:
                    pass

            return '\n'.join(lines).strip()
        except Exception as e:
            logger.error(f"PPT文本提取失败: {e}")
            return f"PPT文本提取失败: {str(e)}"

    @staticmethod
    def extract_text_from_image(file_path: str) -> str:
        try:
            try:
                from PIL import Image
                import pytesseract
            except ImportError:
                return '图片文本提取失败: OCR 依赖 pillow/pytesseract 未安装'

            with Image.open(file_path) as image:
                if image.mode not in ('L', 'RGB'):
                    image = image.convert('RGB')
                try:
                    text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                except Exception:
                    text = pytesseract.image_to_string(image)

            text = text.strip()
            return text or '图片OCR未识别到文本'
        except Exception as e:
            logger.error(f"图片文本提取失败: {e}")
            return f"图片文本提取失败: {str(e)}"

    @classmethod
    def _archive_member_allowed(cls, name: str, size: int, stats: Dict[str, Any], allow_archives: bool = False) -> bool:
        normalized_name = str(name or '').replace('\\', '/')
        if not normalized_name or normalized_name.endswith('/') or normalized_name.startswith('__MACOSX/'):
            return False
        if '..' in normalized_name.split('/'):
            stats['skipped'].append(f"{normalized_name}: unsafe path")
            return False
        if not cls.is_supported_document_file(normalized_name, allow_archives=allow_archives):
            stats['skipped'].append(f"{normalized_name}: unsupported type")
            return False
        if size > cls.MAX_ARCHIVE_MEMBER_BYTES:
            stats['skipped'].append(f"{normalized_name}: too large")
            return False
        if stats['count'] >= cls.MAX_ARCHIVE_FILES:
            stats['skipped'].append(f"{normalized_name}: file limit reached")
            return False
        if stats['total_size'] + max(size, 0) > cls.MAX_ARCHIVE_TOTAL_BYTES:
            stats['skipped'].append(f"{normalized_name}: archive size limit reached")
            return False

        stats['count'] += 1
        stats['total_size'] += max(size, 0)
        return True

    @classmethod
    def _extract_member_stream(cls, stream, member_name: str, depth: int = 0) -> str:
        suffix = cls.get_file_extension(member_name)
        if suffix == '.tar.gz':
            suffix = '.tgz'

        temp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_path = temp_file.name
                shutil.copyfileobj(stream, temp_file)
            return cls.extract_text_from_file(temp_path, filename=member_name, depth=depth + 1)
        finally:
            if temp_path:
                try:
                    os.unlink(temp_path)
                except OSError:
                    pass

    @classmethod
    def _format_archive_parts(cls, parts: list[str], stats: Dict[str, Any]) -> str:
        output = '\n\n'.join(part for part in parts if part.strip()).strip()
        if stats.get('skipped'):
            skipped_preview = '\n'.join(f"- {item}" for item in stats['skipped'][:20])
            skipped_count = len(stats['skipped'])
            suffix = f"\n\n[Archive skipped files: {skipped_count}]\n{skipped_preview}"
            output = f"{output}{suffix}" if output else suffix.strip()
        return cls.truncate_text(output or '压缩包中未找到可解析的需求文档文件')

    @classmethod
    def _extract_text_from_zip(cls, file_path: str, depth: int = 0) -> str:
        stats = {'count': 0, 'total_size': 0, 'skipped': []}
        parts = []
        with zipfile.ZipFile(file_path, 'r') as archive:
            for info in archive.infolist():
                if info.is_dir():
                    continue
                if not cls._archive_member_allowed(info.filename, info.file_size, stats):
                    continue
                with archive.open(info) as stream:
                    text = cls._extract_member_stream(stream, info.filename, depth=depth)
                if text.strip():
                    parts.append(f"## {info.filename}\n{text.strip()}")
        return cls._format_archive_parts(parts, stats)

    @classmethod
    def _extract_text_from_tar(cls, file_path: str, depth: int = 0) -> str:
        stats = {'count': 0, 'total_size': 0, 'skipped': []}
        parts = []
        with tarfile.open(file_path, 'r:*') as archive:
            for member in archive.getmembers():
                if not member.isfile():
                    continue
                if not cls._archive_member_allowed(member.name, member.size, stats):
                    continue
                stream = archive.extractfile(member)
                if stream is None:
                    continue
                with stream:
                    text = cls._extract_member_stream(stream, member.name, depth=depth)
                if text.strip():
                    parts.append(f"## {member.name}\n{text.strip()}")
        return cls._format_archive_parts(parts, stats)

    @classmethod
    def _extract_text_from_rar_with_python(cls, file_path: str, depth: int = 0) -> Optional[str]:
        try:
            import rarfile
        except ImportError:
            return None

        stats = {'count': 0, 'total_size': 0, 'skipped': []}
        parts = []
        try:
            with rarfile.RarFile(file_path, 'r') as archive:
                for info in archive.infolist():
                    if info.isdir():
                        continue
                    if not cls._archive_member_allowed(info.filename, info.file_size, stats):
                        continue
                    with archive.open(info) as stream:
                        text = cls._extract_member_stream(stream, info.filename, depth=depth)
                    if text.strip():
                        parts.append(f"## {info.filename}\n{text.strip()}")
            return cls._format_archive_parts(parts, stats)
        except Exception as e:
            logger.warning(f"RAR python extraction failed: {e}")
            return None

    @classmethod
    def _extract_text_from_extracted_tree(cls, root_dir: str, depth: int = 0) -> str:
        stats = {'count': 0, 'total_size': 0, 'skipped': []}
        parts = []
        for current_root, dir_names, file_names in os.walk(root_dir):
            dir_names[:] = [name for name in dir_names if name != '__MACOSX']
            for file_name in file_names:
                path = os.path.join(current_root, file_name)
                if os.path.islink(path):
                    continue
                relative_name = os.path.relpath(path, root_dir).replace('\\', '/')
                try:
                    size = os.path.getsize(path)
                except OSError:
                    size = 0
                if not cls._archive_member_allowed(relative_name, size, stats):
                    continue
                text = cls.extract_text_from_file(path, filename=relative_name, depth=depth + 1)
                if text.strip():
                    parts.append(f"## {relative_name}\n{text.strip()}")
        return cls._format_archive_parts(parts, stats)

    @classmethod
    def _extract_text_from_archive_with_tool(cls, file_path: str, depth: int = 0) -> str:
        tool_commands = []
        if shutil.which('unar'):
            tool_commands.append(['unar', '-quiet', '-force-overwrite'])
        if shutil.which('bsdtar'):
            tool_commands.append(['bsdtar'])
        if shutil.which('7z'):
            tool_commands.append(['7z'])

        if not tool_commands:
            return 'RAR文本提取失败: 未安装可用的 unar/bsdtar/7z 解压工具'

        last_error = ''
        for tool in tool_commands:
            with tempfile.TemporaryDirectory() as temp_dir:
                if tool[0] == 'unar':
                    command = [*tool, '-output-directory', temp_dir, file_path]
                elif tool[0] == 'bsdtar':
                    command = [*tool, '-xf', file_path, '-C', temp_dir]
                else:
                    command = [*tool, 'x', f'-o{temp_dir}', '-y', file_path]

                result = subprocess.run(command, capture_output=True, text=True, timeout=120, check=False)
                if result.returncode == 0:
                    return cls._extract_text_from_extracted_tree(temp_dir, depth=depth)
                last_error = (result.stderr or result.stdout or '').strip()

        return f"RAR文本提取失败: {last_error or '外部解压工具执行失败'}"

    @classmethod
    def extract_text_from_archive(cls, file_path: str, depth: int = 0, filename: Optional[str] = None) -> str:
        extension = cls.get_file_extension(filename or file_path)
        try:
            if extension == '.zip':
                return cls._extract_text_from_zip(file_path, depth=depth)
            if extension in ('.tar', '.tar.gz', '.tgz'):
                return cls._extract_text_from_tar(file_path, depth=depth)
            if extension == '.rar':
                text = cls._extract_text_from_rar_with_python(file_path, depth=depth)
                if text is not None:
                    return text
                return cls._extract_text_from_archive_with_tool(file_path, depth=depth)
            if zipfile.is_zipfile(file_path):
                return cls._extract_text_from_zip(file_path, depth=depth)
            if tarfile.is_tarfile(file_path):
                return cls._extract_text_from_tar(file_path, depth=depth)
            try:
                import rarfile
                if rarfile.is_rarfile(file_path):
                    text = cls._extract_text_from_rar_with_python(file_path, depth=depth)
                    if text is not None:
                        return text
                    return cls._extract_text_from_archive_with_tool(file_path, depth=depth)
            except ImportError:
                pass
            return '不支持的压缩包类型'
        except Exception as e:
            logger.error(f"压缩包文本提取失败: {e}")
            return f"压缩包文本提取失败: {str(e)}"

    @classmethod
    def extract_text_from_file(
        cls,
        file_path: str,
        filename: Optional[str] = None,
        document_type: Optional[str] = None,
        depth: int = 0,
    ) -> str:
        source_name = filename or file_path
        extension = cls.get_file_extension(source_name)

        if extension == '.pdf' or document_type == 'pdf':
            return cls.extract_text_from_pdf(file_path)
        if extension == '.docx':
            return cls.extract_text_from_docx(file_path)
        if extension == '.doc':
            return cls.extract_text_from_legacy_office(file_path, file_type='Word')
        if extension in ('.txt', '.md') or document_type == 'txt':
            return cls.extract_text_from_txt(file_path)
        if extension == '.xmind' or document_type == 'xmind':
            return cls.extract_text_from_xmind(file_path)
        if extension in ('.png', '.jpg', '.jpeg') or document_type == 'image':
            return cls.extract_text_from_image(file_path)
        if extension in ('.xls', '.xlsx') or document_type == 'excel':
            return cls.extract_text_from_excel(file_path)
        if extension in ('.ppt', '.pptx') or document_type == 'ppt':
            return cls.extract_text_from_ppt(file_path)
        if extension in cls.ARCHIVE_EXTENSIONS or extension == '.gz' or document_type == 'archive':
            if depth >= 2:
                return '压缩包嵌套层级过深，已跳过'
            return cls.extract_text_from_archive(file_path, depth=depth, filename=source_name)
        return "不支持的文档类型"

    @classmethod
    def extract_text(cls, document: RequirementDocument) -> str:
        file_path = document.file.path
        return cls.extract_text_from_file(
            file_path,
            filename=document.file.name,
            document_type=document.document_type,
        )


class AIService:
    """AI服务类 - 模拟大模型调用"""
    
    @staticmethod
    async def analyze_requirements(text: str, document_title: str = "") -> Dict[str, Any]:
        """
        先进的需求分析 - 使用新的智能分析引擎
        
        Args:
            text: 需求文档文本内容
            document_title: 文档标题
            
        Returns:
            Dict包含分析报告、结构化需求等信息
        """
        try:
            # 直接导入并使用先进分析器
            from apps.requirement_analysis.advanced_analyzer import advanced_analyzer
            
            logger.info(f"使用先进分析器分析需求，文档标题: {document_title}")
            
            # 使用先进分析器进行分析
            result = await advanced_analyzer.analyze_requirements_advanced(text, document_title)
            
            # 转换为原系统期望的格式
            analysis_report = result.get("analysis_report", "")
            structured_requirements = result.get("structured_requirements", {})
            requirements_list = structured_requirements.get("requirements", [])
            
            # 计算分析时间（模拟）
            import time
            analysis_time = time.time() % 10 + 2  # 2-12秒之间的模拟时间
            
            logger.info(f"先进需求分析完成，识别需求{len(requirements_list)}个")
            
            return {
                "analysis_report": analysis_report,
                "requirements": requirements_list,
                "requirements_count": len(requirements_list),
                "analysis_time": analysis_time,
                "quality_assessment": result.get("quality_assessment", {}),
                "risk_analysis": result.get("risk_analysis", {})
            }
            
        except Exception as e:
            logger.error(f"先进需求分析失败: {e}")
            logger.info("使用备用分析方法")
            # fallback到原来的分析逻辑
            return await AIService._fallback_analyze_requirements(text, document_title)
    
    @staticmethod
    async def _fallback_analyze_requirements(text: str, document_title: str = "") -> Dict[str, Any]:
        """备用需求分析方法"""
        # 模拟AI分析过程
        await asyncio.sleep(2)
        
        # 这里应该调用真实的大模型API
        # 现在返回改进的模拟数据
        analysis_report = f"""
# 需求分析报告

## 文档概述
基于提供的需求文档"{document_title}"，共识别出以下主要需求模块和功能点。

## 主要功能模块
1. 用户管理模块
2. 数据处理模块  
3. 报告生成模块
4. 系统配置模块

## 详细需求分析
基于文档内容分析，识别出以下具体需求：

### 功能需求
- 用户认证和权限管理
- 数据录入和维护功能
- 业务流程处理
- 报表和统计功能

### 非功能需求
- 系统性能要求：响应时间 < 3秒
- 安全性要求：数据加密存储
- 可用性要求：99.5%系统可用率
- 兼容性要求：支持主流浏览器

## 风险评估
- 技术实现风险：中等
- 进度风险：低
- 资源风险：低

## 建议
建议采用敏捷开发模式，分阶段实施各功能模块。

分析时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
        """
        
        # 生成基础的结构化需求
        requirements = [
            {
                "requirement_id": "REQ-001",
                "requirement_name": "用户认证管理", 
                "requirement_type": "functional",
                "parent_requirement": None,
                "module": "用户管理",
                "requirement_level": "high",
                "reviewer": "admin",
                "estimated_hours": 16,
                "description": "作为一名系统用户，我希望通过用户名和密码登录系统，这样可以确保系统安全性并获得个性化服务。",
                "acceptance_criteria": "用户能够使用有效凭证成功登录系统，无效凭证登录失败，系统记录登录日志。"
            },
            {
                "requirement_id": "REQ-002",
                "requirement_name": "数据管理功能",
                "requirement_type": "functional", 
                "parent_requirement": None,
                "module": "数据管理",
                "requirement_level": "high",
                "reviewer": "admin",
                "estimated_hours": 24,
                "description": "作为一名数据操作员，我希望能够对系统数据进行增删改查操作，这样可以有效管理业务信息。",
                "acceptance_criteria": "数据操作功能正常，数据完整性得到保证，操作权限控制有效。"
            },
            {
                "requirement_id": "REQ-003",
                "requirement_name": "报表统计功能",
                "requirement_type": "functional",
                "parent_requirement": None,
                "module": "报表管理",
                "requirement_level": "medium", 
                "reviewer": "admin",
                "estimated_hours": 20,
                "description": "作为一名管理人员，我希望能够生成各类业务报表和统计图表，这样可以直观了解业务数据和趋势。",
                "acceptance_criteria": "系统能够生成多种格式的报表，数据准确，支持导出功能。"
            }
        ]
        
        return {
            "analysis_report": analysis_report,
            "requirements": requirements,
            "requirements_count": len(requirements)
        }
    
    @staticmethod
    async def generate_test_cases(requirement: BusinessRequirement, test_level: str, test_priority: str, count: int) -> List[Dict[str, Any]]:
        """生成测试用例 - 大模型A"""
        # 模拟AI生成过程
        await asyncio.sleep(1)
        
        # 生成唯一case_id的辅助函数
        def generate_unique_case_id(req, base_index):
            """生成唯一的测试用例ID"""
            base_case_id = f"TC-{req.requirement_id}-{base_index:03d}"
            case_id = base_case_id
            counter = 1
            
            # 检查是否已存在，如果存在则添加后缀
            from .models import GeneratedTestCase
            while GeneratedTestCase.objects.filter(requirement=req, case_id=case_id).exists():
                case_id = f"{base_case_id}_{counter}"
                counter += 1
            
            return case_id
        
        # 获取该需求现有测试用例的数量，作为起始索引
        from .models import GeneratedTestCase
        existing_count = GeneratedTestCase.objects.filter(requirement=requirement).count()
        
        # 根据需求生成测试用例
        test_cases = []
        for i in range(count):
            case_id = generate_unique_case_id(requirement, existing_count + i + 1)
            
            # 根据需求类型生成不同的测试用例
            if "登录" in requirement.requirement_name:
                test_cases.append({
                    "case_id": case_id,
                    "title": f"验证用户使用有效凭证登录系统的认证流程和权限获取",
                    "priority": test_priority,
                    "precondition": "系统正常运行，测试用户账号已创建",
                    "test_steps": "1. 打开登录页面\n2. 输入有效的用户名和密码\n3. 点击登录按钮\n4. 检查登录结果和页面跳转",
                    "expected_result": "用户成功登录系统，跳转到主页面，显示用户信息和相应权限功能"
                })
            elif "数据" in requirement.requirement_name:
                test_cases.append({
                    "case_id": case_id,
                    "title": f"测试数据录入功能在各种输入场景下的验证机制和保存结果",
                    "priority": test_priority,
                    "precondition": "系统正常运行，用户已登录具备数据操作权限",
                    "test_steps": "1. 进入数据录入页面\n2. 填写必填字段信息\n3. 提交数据\n4. 验证数据保存结果",
                    "expected_result": "数据成功保存到数据库，页面显示保存成功提示，可以查询到新录入的数据"
                })
            elif "报告" in requirement.requirement_name:
                test_cases.append({
                    "case_id": case_id,
                    "title": f"验证报告生成功能在不同格式和数据量下的处理能力和输出质量",
                    "priority": test_priority, 
                    "precondition": "系统正常运行，存在可用于生成报告的数据",
                    "test_steps": "1. 进入报告生成页面\n2. 选择报告类型和参数\n3. 点击生成报告\n4. 检查生成的报告内容和格式",
                    "expected_result": "报告成功生成，内容准确完整，格式符合要求，可以正常下载"
                })
            else:
                test_cases.append({
                    "case_id": case_id,
                    "title": f"验证{requirement.requirement_name}功能的基本操作流程和预期结果",
                    "priority": test_priority,
                    "precondition": "系统正常运行，用户已登录",
                    "test_steps": f"1. 访问{requirement.requirement_name}功能\n2. 执行主要操作步骤\n3. 验证操作结果",
                    "expected_result": f"{requirement.requirement_name}功能正常工作，操作结果符合预期"
                })
        
        return test_cases
    
    @staticmethod
    async def review_test_cases(test_cases: List[GeneratedTestCase], review_criteria: str) -> Dict[str, Any]:
        """评审测试用例 - 大模型B"""
        # 模拟AI评审过程
        await asyncio.sleep(1.5)
        
        reviewed_cases = []
        for test_case in test_cases:
            # 模拟评审逻辑
            review_score = 85  # 模拟评分
            
            review_comments = f"""
评审意见:
1. 测试用例标题清晰明确，能够准确描述测试目的
2. 测试步骤详细具体，具有良好的可执行性
3. 预期结果明确，便于验证
4. 建议补充异常场景的测试覆盖

评审分数: {review_score}/100
评审状态: 通过
"""
            
            reviewed_cases.append({
                "test_case_id": test_case.id,
                "review_score": review_score,
                "review_comments": review_comments.strip(),
                "status": "reviewed" if review_score >= 80 else "rejected"
            })
        
        return {
            "reviewed_cases": reviewed_cases,
            "overall_score": sum(case["review_score"] for case in reviewed_cases) / len(reviewed_cases),
            "pass_rate": len([case for case in reviewed_cases if case["status"] == "reviewed"]) / len(reviewed_cases) * 100
        }


class RequirementAnalysisService:
    """需求分析服务"""
    
    @classmethod
    def create_analysis_task(cls, document: RequirementDocument, task_type: str) -> AnalysisTask:
        """创建分析任务"""
        task_id = f"{task_type}_{uuid.uuid4().hex[:8]}"
        
        task = AnalysisTask.objects.create(
            task_id=task_id,
            task_type=task_type,
            document=document,
            status='pending'
        )
        
        return task
    
    @classmethod
    async def process_document_analysis(cls, document: RequirementDocument) -> RequirementAnalysis:
        """处理文档分析"""
        # 创建分析任务
        task = cls.create_analysis_task(document, 'requirement_analysis')
        
        try:
            # 更新任务状态
            task.status = 'running'
            task.started_at = datetime.now()
            task.progress = 10
            task.save()
            
            # 提取文档文本
            if not document.extracted_text:
                document.extracted_text = DocumentProcessor.extract_text(document)
                document.save()
            
            task.progress = 30
            task.save()
            
            # 调用AI分析
            start_time = time.time()
            analysis_result = await AIService.analyze_requirements(
                document.extracted_text, 
                document.title
            )
            analysis_time = time.time() - start_time
            
            task.progress = 70
            task.save()
            
            # 创建分析记录
            analysis = RequirementAnalysis.objects.create(
                document=document,
                analysis_report=analysis_result['analysis_report'],
                requirements_count=analysis_result['requirements_count'],
                analysis_time=analysis_time
            )
            
            # 保存需求数据
            for req_data in analysis_result['requirements']:
                BusinessRequirement.objects.create(
                    analysis=analysis,
                    **req_data
                )
            
            # 更新文档状态
            document.status = 'analyzed'
            document.save()
            
            # 完成任务
            task.status = 'completed'
            task.completed_at = datetime.now()
            task.progress = 100
            task.result = analysis_result
            task.save()
            
            return analysis
            
        except Exception as e:
            logger.error(f"文档分析失败: {e}")
            
            # 更新任务状态
            task.status = 'failed'
            task.error_message = str(e)
            task.completed_at = datetime.now()
            task.save()
            
            # 更新文档状态
            document.status = 'failed'
            document.save()
            
            raise e
    
    @classmethod
    async def generate_test_cases_for_requirements(cls, requirement_ids: List[int], test_level: str, test_priority: str, test_case_count: int) -> List[GeneratedTestCase]:
        """为需求生成测试用例"""
        generated_cases = []
        
        for req_id in requirement_ids:
            try:
                requirement = BusinessRequirement.objects.get(id=req_id)
                
                # 调用AI生成测试用例
                test_cases_data = await AIService.generate_test_cases(
                    requirement, test_level, test_priority, test_case_count
                )
                
                # 保存生成的测试用例
                for case_data in test_cases_data:
                    test_case = GeneratedTestCase.objects.create(
                        requirement=requirement,
                        case_id=case_data['case_id'],
                        title=case_data['title'],
                        priority=case_data['priority'],
                        precondition=case_data['precondition'],
                        test_steps=case_data['test_steps'],
                        expected_result=case_data['expected_result'],
                        generated_by_ai='AI-A'
                    )
                    generated_cases.append(test_case)
                    
            except BusinessRequirement.DoesNotExist:
                logger.error(f"需求ID {req_id} 不存在")
                continue
            except Exception as e:
                logger.error(f"为需求 {req_id} 生成测试用例失败: {e}")
                continue
        
        return generated_cases
    
    @classmethod
    async def review_test_cases(cls, test_case_ids: List[int], review_criteria: str) -> Dict[str, Any]:
        """评审测试用例"""
        test_cases = GeneratedTestCase.objects.filter(id__in=test_case_ids)
        
        # 调用AI评审
        review_result = await AIService.review_test_cases(list(test_cases), review_criteria)
        
        # 更新测试用例状态
        for case_review in review_result['reviewed_cases']:
            try:
                test_case = GeneratedTestCase.objects.get(id=case_review['test_case_id'])
                test_case.status = case_review['status']
                test_case.review_comments = case_review['review_comments']
                test_case.reviewed_by_ai = 'AI-B'
                test_case.save()
            except GeneratedTestCase.DoesNotExist:
                logger.error(f"测试用例ID {case_review['test_case_id']} 不存在")
                continue
        
        return review_result
